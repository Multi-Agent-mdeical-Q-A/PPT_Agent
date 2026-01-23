#!/usr/bin/env python3
"""
batch_ppt_to_marp.py (v2.0)
===========================
批量将 PPTX 课件转换为 Marp Markdown 格式，专为多轮 AI 内容重构优化。

核心改进 (v2.0):
- 🎯 保留原始标题：无标题时智能回退，绝不用页码占位
- 📊 丰富的结构化元数据：层级关系、内容类型、布局信息
- 🔗 幻灯片间关系推断：章节检测、内容连续性标记
- 📝 增强的文本提取：保留层级缩进、识别列表 vs 段落
- 🖼️ 智能图片处理：位置信息、尺寸比例、多图布局建议
- 🏷️ AI 友好的标签系统：方便下游工具解析和处理
- 📈 内容质量指标：文字密度、图文比例等

使用方法：
    python batch_ppt_to_marp.py [--verbose] [--output-dir DIR]

依赖：
    pip install python-pptx --break-system-packages
"""

import os
import sys
import re
import json
import hashlib
from pathlib import Path
from typing import Optional, Tuple, List, Dict, Any
from dataclasses import dataclass, field, asdict
from datetime import datetime
from enum import Enum

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("=" * 60)
    print("Error: python-pptx 未安装")
    print("请运行以下命令安装:")
    print("  pip install python-pptx --break-system-packages")
    print("=" * 60)
    sys.exit(1)


# =============================================================================
# 数据结构定义
# =============================================================================

class SlideType(Enum):
    """幻灯片类型分类"""
    TITLE = "title"           # 标题页/封面
    SECTION = "section"       # 章节分隔页
    CONTENT = "content"       # 普通内容页
    IMAGE_ONLY = "image_only" # 纯图片页
    TEXT_ONLY = "text_only"   # 纯文字页
    MIXED = "mixed"           # 图文混排
    BLANK = "blank"           # 空白页
    ENDING = "ending"         # 结束页


class ContentDensity(Enum):
    """内容密度评估"""
    SPARSE = "sparse"     # 稀疏（适合扩展）
    NORMAL = "normal"     # 正常
    DENSE = "dense"       # 密集（可能需要拆分）


@dataclass
class ImageInfo:
    """图片信息"""
    filename: str
    original_name: str
    width_emu: int
    height_emu: int
    left_emu: int
    top_emu: int
    aspect_ratio: float
    position_hint: str  # "left", "right", "center", "full"
    size_hint: str      # "small", "medium", "large", "full"
    
    def to_marp_directive(self, asset_path: str) -> str:
        """生成 Marp 图片指令"""
        # 根据位置和大小生成最佳布局
        if self.size_hint == "full":
            return f"![bg contain]({asset_path})"
        elif self.position_hint == "right":
            width_pct = "40%" if self.size_hint == "large" else "35%"
            return f"![bg right:{width_pct} fit]({asset_path})"
        elif self.position_hint == "left":
            width_pct = "40%" if self.size_hint == "large" else "35%"
            return f"![bg left:{width_pct} fit]({asset_path})"
        else:
            return f"![bg right:35% fit]({asset_path})"


@dataclass
class TextBlock:
    """文本块信息"""
    text: str
    level: int              # 缩进层级 (0=顶级)
    is_title: bool
    is_subtitle: bool
    is_bullet: bool
    font_size: Optional[float]  # 字号（磅）
    is_bold: bool
    shape_type: str         # "title", "body", "textbox", "other"


@dataclass
class SlideData:
    """单张幻灯片的结构化数据"""
    index: int
    total: int
    
    # 内容
    title: Optional[str]
    subtitle: Optional[str]
    text_blocks: List[TextBlock]
    images: List[ImageInfo]
    speaker_notes: str
    
    # 元数据
    slide_type: SlideType
    content_density: ContentDensity
    has_animation: bool
    layout_name: str
    
    # AI 辅助信息
    is_section_start: bool
    section_title: Optional[str]
    estimated_speak_time_sec: int
    key_terms: List[str]
    
    # 关系
    continues_from_previous: bool
    continues_to_next: bool


# =============================================================================
# 工具函数
# =============================================================================

def sanitize_filename(name: str) -> str:
    """清理文件名"""
    sanitized = re.sub(r'[<>:"/\\|?*]', '_', name)
    sanitized = re.sub(r'\s+', '_', sanitized)
    sanitized = re.sub(r'_+', '_', sanitized)
    return sanitized.strip('_')


def emu_to_inches(emu: int) -> float:
    """EMU 转英寸"""
    return emu / 914400


def estimate_speak_time(text_content: str, notes: str) -> int:
    """估算演讲时间（秒）"""
    # 中文约 150 字/分钟，英文约 130 词/分钟
    total_chars = len(text_content) + len(notes)
    # 粗略估算：每个字符约 0.4 秒
    return max(30, int(total_chars * 0.4))


def extract_key_terms(text: str) -> List[str]:
    """提取关键术语（简单实现）"""
    # 提取引号内容、大写缩写、专有名词等
    terms = []
    
    # 引号内容
    quoted = re.findall(r'[""「」『』]([^""「」『』]+)[""「」『』]', text)
    terms.extend(quoted)
    
    # 英文缩写 (2-5个大写字母)
    abbrevs = re.findall(r'\b[A-Z]{2,5}\b', text)
    terms.extend(abbrevs)
    
    # 去重并限制数量
    seen = set()
    unique_terms = []
    for t in terms:
        t_lower = t.lower()
        if t_lower not in seen and len(t) > 1:
            seen.add(t_lower)
            unique_terms.append(t)
    
    return unique_terms[:10]


def detect_slide_type(title: Optional[str], text_blocks: List[TextBlock], 
                      images: List[ImageInfo], layout_name: str) -> SlideType:
    """检测幻灯片类型"""
    has_text = bool(text_blocks)
    has_images = bool(images)
    
    # 检查布局名称中的关键词
    layout_lower = layout_name.lower()
    if any(kw in layout_lower for kw in ['title', '标题', '封面']):
        if 'section' in layout_lower or '节' in layout_lower:
            return SlideType.SECTION
        return SlideType.TITLE
    
    if any(kw in layout_lower for kw in ['blank', '空白']):
        return SlideType.BLANK
    
    if any(kw in layout_lower for kw in ['end', '结束', 'thank', '谢谢']):
        return SlideType.ENDING
    
    # 基于内容判断
    if not has_text and not has_images:
        return SlideType.BLANK
    elif has_images and not has_text:
        return SlideType.IMAGE_ONLY
    elif has_text and not has_images:
        return SlideType.TEXT_ONLY
    else:
        return SlideType.MIXED


def assess_content_density(text_blocks: List[TextBlock], images: List[ImageInfo]) -> ContentDensity:
    """评估内容密度"""
    total_text_len = sum(len(tb.text) for tb in text_blocks)
    num_items = len(text_blocks) + len(images)
    
    if total_text_len < 50 and num_items <= 2:
        return ContentDensity.SPARSE
    elif total_text_len > 500 or num_items > 8:
        return ContentDensity.DENSE
    else:
        return ContentDensity.NORMAL


def is_continuation_title(title: str) -> bool:
    """检查是否是延续性标题（如 "xxx（续）"）"""
    if not title:
        return False
    patterns = [
        r'[\(（][续继][\)）]',
        r'cont[\'.]?d?',
        r'continued',
        r'part\s*\d+',
        r'[（\(]\d+[）\)]$',
    ]
    title_lower = title.lower()
    return any(re.search(p, title_lower) for p in patterns)


def detect_section_start(title: str, layout_name: str, prev_title: Optional[str]) -> Tuple[bool, Optional[str]]:
    """检测是否是章节开始"""
    if not title:
        return False, None
    
    layout_lower = layout_name.lower()
    
    # 布局名称包含 section
    if 'section' in layout_lower or '节' in layout_lower:
        return True, title
    
    # 标题格式检测：数字开头、"第X章/节" 等
    section_patterns = [
        r'^第[一二三四五六七八九十\d]+[章节讲课]',
        r'^\d+[\.\、]\s*\S',
        r'^[IVX]+[\.\、]\s*\S',
        r'^(Chapter|Lecture|Week|Part)\s*\d+',
    ]
    
    for pattern in section_patterns:
        if re.match(pattern, title, re.IGNORECASE):
            return True, title
    
    return False, None


# =============================================================================
# 核心提取函数
# =============================================================================

def get_image_position_hint(left_emu: int, top_emu: int, width_emu: int, 
                            slide_width: int, slide_height: int) -> str:
    """根据位置判断图片应该放在哪里"""
    center_x = left_emu + width_emu / 2
    slide_center = slide_width / 2
    
    # 判断水平位置
    if center_x < slide_center * 0.6:
        return "left"
    elif center_x > slide_center * 1.4:
        return "right"
    else:
        return "center"


def get_image_size_hint(width_emu: int, height_emu: int, 
                        slide_width: int, slide_height: int) -> str:
    """根据尺寸判断图片大小级别"""
    area_ratio = (width_emu * height_emu) / (slide_width * slide_height)
    
    if area_ratio > 0.5:
        return "full"
    elif area_ratio > 0.25:
        return "large"
    elif area_ratio > 0.1:
        return "medium"
    else:
        return "small"


def extract_images_enhanced(slide, slide_num: int, output_dir: Path,
                            slide_width: int, slide_height: int) -> List[ImageInfo]:
    """增强版图片提取"""
    images = []
    img_counter = 1
    
    def process_shape(shape):
        nonlocal img_counter
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext.lstrip('.')
                
                # 获取位置和尺寸
                left = getattr(shape, 'left', 0)
                top = getattr(shape, 'top', 0)
                width = getattr(shape, 'width', 0)
                height = getattr(shape, 'height', 0)
                
                # 计算宽高比
                aspect = width / height if height > 0 else 1.0
                
                # 生成文件名
                img_filename = f"slide_{slide_num:02d}_img_{img_counter:02d}.{ext}"
                img_path = output_dir / img_filename
                
                # 保存图片
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                
                # 尝试获取原始文件名
                orig_name = getattr(image, 'filename', img_filename)
                
                info = ImageInfo(
                    filename=img_filename,
                    original_name=orig_name,
                    width_emu=width,
                    height_emu=height,
                    left_emu=left,
                    top_emu=top,
                    aspect_ratio=round(aspect, 2),
                    position_hint=get_image_position_hint(left, top, width, slide_width, slide_height),
                    size_hint=get_image_size_hint(width, height, slide_width, slide_height)
                )
                
                images.append(info)
                img_counter += 1
                
        except Exception:
            pass
    
    for shape in slide.shapes:
        process_shape(shape)
        # 处理组合形状
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                process_shape(sub_shape)
    
    # 按位置排序（从左到右，从上到下）
    images.sort(key=lambda x: (x.top_emu, x.left_emu))
    
    return images


def extract_text_enhanced(slide) -> Tuple[Optional[str], Optional[str], List[TextBlock]]:
    """
    增强版文本提取
    
    Returns:
        (标题, 副标题, 文本块列表)
    """
    title = None
    subtitle = None
    text_blocks = []
    
    # 用于追踪已处理的占位符
    title_found = False
    subtitle_found = False
    
    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            shape_type = "other"
            is_title_shape = False
            is_subtitle_shape = False
            
            # 检查占位符类型
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    # TITLE=1, CENTER_TITLE=3, SUBTITLE=4, BODY=2
                    if ph_type in [1, 3]:
                        is_title_shape = True
                        shape_type = "title"
                    elif ph_type == 4:
                        is_subtitle_shape = True
                        shape_type = "subtitle"
                    elif ph_type == 2:
                        shape_type = "body"
                except:
                    pass
            
            # 提取段落
            for para_idx, para in enumerate(text_frame.paragraphs):
                # 收集段落文本
                para_text = ""
                is_bold = False
                font_size = None
                
                for run in para.runs:
                    para_text += run.text
                    # 获取格式信息
                    if run.font.bold:
                        is_bold = True
                    if run.font.size:
                        font_size = run.font.size.pt
                
                para_text = para_text.strip()
                if not para_text:
                    continue
                
                # 获取缩进层级
                level = para.level if para.level is not None else 0
                
                # 检查是否是列表项
                is_bullet = level > 0 or (hasattr(para, 'bullet') and para.bullet is not None)
                
                # 处理标题
                if is_title_shape and not title_found:
                    title = para_text
                    title_found = True
                    continue
                
                # 处理副标题
                if is_subtitle_shape and not subtitle_found:
                    subtitle = para_text
                    subtitle_found = True
                    continue
                
                # 创建文本块
                block = TextBlock(
                    text=para_text,
                    level=level,
                    is_title=is_title_shape,
                    is_subtitle=is_subtitle_shape,
                    is_bullet=is_bullet,
                    font_size=font_size,
                    is_bold=is_bold,
                    shape_type=shape_type
                )
                text_blocks.append(block)
                
        except Exception:
            continue
    
    return title, subtitle, text_blocks


def extract_speaker_notes(slide) -> str:
    """提取演讲者备注"""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes = []
            for para in notes_text_frame.paragraphs:
                para_text = "".join(run.text for run in para.runs).strip()
                if para_text:
                    notes.append(para_text)
            return "\n".join(notes)
    except Exception:
        pass
    return ""


def get_layout_name(slide) -> str:
    """获取幻灯片布局名称"""
    try:
        return slide.slide_layout.name
    except:
        return "Unknown"


# =============================================================================
# Markdown 生成
# =============================================================================

def generate_marp_header(title: str, total_slides: int, source_file: str) -> str:
    """生成增强的 Marp YAML 头部"""
    timestamp = datetime.now().isoformat()
    return f"""---
marp: true
theme: default
paginate: true
size: 16:9
header: "{title}"
footer: "Course Refactoring Draft"
# === AI REFACTORING METADATA ===
# source_file: {source_file}
# total_slides: {total_slides}
# extracted_at: {timestamp}
# version: 2.0
---

"""


def generate_slide_markdown(slide_data: SlideData, asset_folder: str) -> str:
    """生成单张幻灯片的 Markdown"""
    lines = []
    
    # === 标题 ===
    # 优先使用原始标题，绝不使用页码
    if slide_data.title:
        lines.append(f"# {slide_data.title}")
    elif slide_data.subtitle:
        # 回退到副标题
        lines.append(f"# {slide_data.subtitle}")
    elif slide_data.text_blocks:
        # 尝试从第一个文本块推断标题
        first_block = slide_data.text_blocks[0]
        if first_block.is_bold or (first_block.font_size and first_block.font_size > 20):
            lines.append(f"# {first_block.text}")
        else:
            # 使用描述性占位符，方便 AI 后续处理
            lines.append("# (Untitled - needs AI review)")
    elif slide_data.images:
        # 纯图片页
        lines.append("# (Visual Content)")
    else:
        lines.append("# (Untitled - needs AI review)")
    
    lines.append("")
    
    # === 图片 ===
    if slide_data.images:
        primary_img = slide_data.images[0]
        asset_path = f"assets/{asset_folder}/{primary_img.filename}"
        lines.append(primary_img.to_marp_directive(asset_path))
        lines.append("")
        
        # 额外图片记录在注释中
        if len(slide_data.images) > 1:
            lines.append("<!-- [ADDITIONAL_IMAGES]")
            for img in slide_data.images[1:]:
                img_path = f"assets/{asset_folder}/{img.filename}"
                lines.append(f"  - {img_path} (position: {img.position_hint}, size: {img.size_hint})")
            lines.append("[END_IMAGES] -->")
            lines.append("")
    
    # === 正文内容 ===
    content_added = False
    if slide_data.text_blocks:
        for block in slide_data.text_blocks:
            # 跳过已用作标题的内容
            if block.text == slide_data.title or block.text == slide_data.subtitle:
                continue
            
            # 根据层级生成缩进
            indent = "  " * block.level
            prefix = "-"
            
            # 清理文本
            clean_text = " ".join(block.text.split())
            if clean_text:
                lines.append(f"{indent}{prefix} {clean_text}")
                content_added = True
        
        if content_added:
            lines.append("")
    
    # === 结构化元数据注释块 ===
    lines.append("<!--")
    lines.append(f"[SLIDE_META]")
    lines.append(f"  position: {slide_data.index}/{slide_data.total}")
    lines.append(f"  type: {slide_data.slide_type.value}")
    lines.append(f"  layout: {slide_data.layout_name}")
    lines.append(f"  density: {slide_data.content_density.value}")
    lines.append(f"  est_time_sec: {slide_data.estimated_speak_time_sec}")
    
    if slide_data.is_section_start:
        lines.append(f"  section_start: true")
        if slide_data.section_title:
            lines.append(f"  section_title: {slide_data.section_title}")
    
    if slide_data.continues_from_previous:
        lines.append(f"  continues_from_previous: true")
    if slide_data.continues_to_next:
        lines.append(f"  continues_to_next: true")
    
    if slide_data.key_terms:
        lines.append(f"  key_terms: {', '.join(slide_data.key_terms)}")
    
    lines.append(f"[END_META]")
    lines.append("")
    
    # === 演讲者备注 ===
    lines.append("[SPEAKER_NOTES]")
    if slide_data.speaker_notes:
        lines.append(slide_data.speaker_notes)
    else:
        lines.append("(No speaker notes)")
    lines.append("[END_NOTES]")
    lines.append("-->")
    
    return "\n".join(lines)


# =============================================================================
# 主处理函数
# =============================================================================

def process_single_pptx(pptx_path: Path, assets_base_dir: Path, 
                        verbose: bool = False) -> Tuple[bool, str, Dict[str, Any]]:
    """
    处理单个 PPTX 文件
    
    Returns:
        (成功标志, 消息, 统计信息)
    """
    stats = {
        "slides": 0,
        "images": 0,
        "notes_count": 0,
        "sections": 0,
        "warnings": []
    }
    
    try:
        stem = pptx_path.stem
        sanitized_stem = sanitize_filename(stem)
        
        # 创建资源目录
        asset_dir = assets_base_dir / sanitized_stem
        asset_dir.mkdir(parents=True, exist_ok=True)
        
        # 输出文件
        md_path = pptx_path.parent / f"{sanitized_stem}.md"
        
        # 打开 PPT
        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)
        stats["slides"] = total_slides
        
        # 获取幻灯片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # === 第一遍：提取所有数据 ===
        slides_data: List[SlideData] = []
        prev_title = None
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            # 提取内容
            images = extract_images_enhanced(slide, slide_idx, asset_dir, slide_width, slide_height)
            title, subtitle, text_blocks = extract_text_enhanced(slide)
            notes = extract_speaker_notes(slide)
            layout_name = get_layout_name(slide)
            
            # 统计
            stats["images"] += len(images)
            if notes:
                stats["notes_count"] += 1
            
            # 检测章节
            is_section, section_title = detect_section_start(title or "", layout_name, prev_title)
            if is_section:
                stats["sections"] += 1
            
            # 检测延续性
            continues_from = is_continuation_title(title) if title else False
            
            # 构建全文用于分析
            all_text = " ".join([title or "", subtitle or ""] + [tb.text for tb in text_blocks] + [notes])
            
            # 创建数据对象
            slide_data = SlideData(
                index=slide_idx,
                total=total_slides,
                title=title,
                subtitle=subtitle,
                text_blocks=text_blocks,
                images=images,
                speaker_notes=notes,
                slide_type=detect_slide_type(title, text_blocks, images, layout_name),
                content_density=assess_content_density(text_blocks, images),
                has_animation=False,  # python-pptx 不直接支持动画检测
                layout_name=layout_name,
                is_section_start=is_section,
                section_title=section_title,
                estimated_speak_time_sec=estimate_speak_time(all_text, notes),
                key_terms=extract_key_terms(all_text),
                continues_from_previous=continues_from,
                continues_to_next=False  # 后面处理
            )
            
            slides_data.append(slide_data)
            prev_title = title
        
        # === 第二遍：标记延续关系 ===
        for i in range(len(slides_data) - 1):
            if slides_data[i + 1].continues_from_previous:
                slides_data[i].continues_to_next = True
        
        # === 生成 Markdown ===
        md_content = generate_marp_header(stem, total_slides, pptx_path.name)
        
        for i, sd in enumerate(slides_data):
            md_content += generate_slide_markdown(sd, sanitized_stem)
            
            # 幻灯片分隔符
            if i < len(slides_data) - 1:
                md_content += "\n\n---\n\n"
            else:
                md_content += "\n"
        
        # 写入文件
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        
        # === 生成配套的 JSON 元数据（方便程序化处理）===
        meta_path = pptx_path.parent / f"{sanitized_stem}_meta.json"
        meta = {
            "source": pptx_path.name,
            "output": md_path.name,
            "extracted_at": datetime.now().isoformat(),
            "stats": stats,
            "structure": {
                "total_slides": total_slides,
                "sections": [
                    {"index": sd.index, "title": sd.section_title}
                    for sd in slides_data if sd.is_section_start
                ]
            },
            "slides": [
                {
                    "index": sd.index,
                    "title": sd.title,
                    "subtitle": sd.subtitle,
                    "type": sd.slide_type.value,
                    "density": sd.content_density.value,
                    "layout": sd.layout_name,
                    "has_notes": bool(sd.speaker_notes),
                    "image_count": len(sd.images),
                    "text_block_count": len(sd.text_blocks),
                    "is_section_start": sd.is_section_start,
                    "section_title": sd.section_title,
                    "key_terms": sd.key_terms,
                    "est_time_sec": sd.estimated_speak_time_sec,
                    "continues_from_previous": sd.continues_from_previous,
                    "continues_to_next": sd.continues_to_next
                }
                for sd in slides_data
            ]
        }
        with open(meta_path, 'w', encoding='utf-8') as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)
        
        msg = f"✓ {total_slides} slides, {stats['images']} imgs, {stats['notes_count']} notes, {stats['sections']} sections"
        return True, msg, stats
        
    except Exception as e:
        stats["warnings"].append(str(e))
        return False, f"✗ Error: {str(e)}", stats


def main():
    """主函数"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="批量将 PPTX 转换为 AI 友好的 Marp Markdown"
    )
    parser.add_argument("--verbose", "-v", action="store_true", help="显示详细输出")
    parser.add_argument("--output-dir", "-o", type=str, help="指定输出目录（默认为当前目录）")
    args = parser.parse_args()
    
    print()
    print("╔" + "═" * 58 + "╗")
    print("║" + "  PPTX → Marp Converter v2.0".center(58) + "║")
    print("║" + "  Optimized for AI-Assisted Refactoring".center(58) + "║")
    print("╚" + "═" * 58 + "╝")
    print()
    
    # 确定工作目录
    current_dir = Path(args.output_dir) if args.output_dir else Path.cwd()
    print(f"📁 Working Directory: {current_dir}")
    print()
    
    # 扫描 PPTX 文件
    pptx_files = sorted(current_dir.glob("*.pptx"))
    
    # 排除临时文件
    pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]
    
    if not pptx_files:
        print("⚠️  No .pptx files found in current directory.")
        print("   Please run this script in the folder containing your PPT files.")
        sys.exit(0)
    
    print(f"📊 Found {len(pptx_files)} PPTX file(s):")
    for f in pptx_files:
        print(f"   • {f.name}")
    print()
    
    # 创建 assets 目录
    assets_dir = current_dir / "assets"
    assets_dir.mkdir(exist_ok=True)
    
    # 处理统计
    results = {
        "success": 0,
        "failed": 0,
        "total_slides": 0,
        "total_images": 0,
        "total_notes": 0
    }
    
    print("─" * 60)
    print("Processing...")
    print("─" * 60)
    
    for idx, pptx_file in enumerate(pptx_files, start=1):
        prefix = f"[{idx:2d}/{len(pptx_files)}]"
        print(f"{prefix} {pptx_file.name}")
        
        success, message, stats = process_single_pptx(pptx_file, assets_dir, args.verbose)
        
        print(f"       {message}")
        
        if success:
            results["success"] += 1
            results["total_slides"] += stats["slides"]
            results["total_images"] += stats["images"]
            results["total_notes"] += stats["notes_count"]
        else:
            results["failed"] += 1
        
        if args.verbose and stats.get("warnings"):
            for warn in stats["warnings"]:
                print(f"       ⚠️  {warn}")
    
    # 最终报告
    print()
    print("─" * 60)
    print("📈 Summary")
    print("─" * 60)
    print(f"   Files processed:  {len(pptx_files)}")
    print(f"   Successful:       {results['success']}")
    print(f"   Failed:           {results['failed']}")
    print()
    print(f"   Total slides:     {results['total_slides']}")
    print(f"   Total images:     {results['total_images']}")
    print(f"   Slides w/ notes:  {results['total_notes']}")
    print()
    print("📂 Output Structure:")
    print(f"   Markdown:  {current_dir}/*.md")
    print(f"   Metadata:  {current_dir}/*_meta.json")
    print(f"   Images:    {assets_dir}/<filename>/")
    print()
    print("🚀 Ready for multi-stage AI refactoring!")
    print()


if __name__ == "__main__":
    main()
